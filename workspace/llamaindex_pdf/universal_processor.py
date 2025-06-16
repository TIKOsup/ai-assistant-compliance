# universal_processor.py
import os
import sys
import pandas as pd
from docx import Document as DocxDocument
import pymupdf4llm
import pytesseract
from pdf2image import convert_from_path
import xml.etree.ElementTree as ET
from datetime import datetime
import warnings
from PIL import Image

# Отключаем предупреждения
warnings.filterwarnings("ignore")
Image.MAX_IMAGE_PIXELS = None

class UniversalDocumentProcessor:
    def __init__(self, input_dir="./regulations", output_dir="./processed_regulations"):
        self.input_dir = input_dir
        self.output_dir = output_dir
        self.processed_files = {}
        
        # Создаем выходную папку
        os.makedirs(output_dir, exist_ok=True)
        
    def extract_text_from_pdf(self, pdf_path):
        """Извлечение текста из PDF с OCR поддержкой"""
        print(f"📄 Обработка PDF: {os.path.basename(pdf_path)}")
        
        try:
            # Сначала пробуем обычное извлечение
            text = pymupdf4llm.to_markdown(pdf_path)
            if len(text.strip()) > 100:
                print(f"✅ Обычное извлечение: {len(text)} символов")
                return text
            else:
                print("⚠️ Мало текста, используем OCR...")
        except Exception as e:
            print(f"⚠️ Ошибка обычного извлечения: {e}")
            print("🔄 Переходим к OCR...")
        
        # OCR обработка
        try:
            print("📷 Конвертация в изображения...")
            images = convert_from_path(pdf_path, dpi=300)
            
            all_text = []
            for i, image in enumerate(images):
                print(f"🔤 OCR страница {i+1}/{len(images)}...")
                
                # Настройки OCR для русского и английского
                custom_config = r'--oem 3 --psm 6 -l rus+eng'
                text = pytesseract.image_to_string(image, config=custom_config)
                
                if text.strip():
                    all_text.append(f"=== Страница {i+1} ===\n{text}")
            
            combined_text = "\n\n".join(all_text)
            print(f"✅ OCR завершен: {len(combined_text)} символов")
            return combined_text
            
        except Exception as e:
            print(f"❌ Ошибка OCR: {e}")
            return None

    def extract_text_from_docx(self, docx_path):
        """Извлечение текста из DOCX файлов"""
        print(f"📄 Обработка DOCX: {os.path.basename(docx_path)}")
        
        try:
            doc = DocxDocument(docx_path)
            text = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text.append(paragraph.text)
            
            # Обработка таблиц
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text.append(" | ".join(row_text))
            
            combined_text = "\n".join(text)
            print(f"✅ DOCX обработан: {len(combined_text)} символов")
            return combined_text
            
        except Exception as e:
            print(f"❌ Ошибка DOCX: {e}")
            return None

    def extract_text_from_doc(self, doc_path):
        """Извлечение текста из старых DOC файлов"""
        print(f"📄 Обработка DOC: {os.path.basename(doc_path)}")
        
        try:
            # Пробуем использовать python-docx (иногда работает с .doc)
            doc = DocxDocument(doc_path)
            text = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text.append(paragraph.text)
            
            combined_text = "\n".join(text)
            if len(combined_text) > 50:
                print(f"✅ DOC обработан: {len(combined_text)} символов")
                return combined_text
            else:
                print("⚠️ Мало текста из DOC файла")
                return f"[DOC файл {os.path.basename(doc_path)} требует дополнительной обработки]"
                
        except Exception as e:
            print(f"⚠️ Не удалось обработать DOC файл: {e}")
            return f"[DOC файл {os.path.basename(doc_path)} не удалось обработать - {str(e)}]"

    def extract_text_from_xlsx(self, xlsx_path):
        """Извлечение данных из Excel файлов"""
        print(f"📄 Обработка XLSX: {os.path.basename(xlsx_path)}")
        
        try:
            # Читаем все листы
            xlsx_file = pd.ExcelFile(xlsx_path)
            all_text = []
            
            for sheet_name in xlsx_file.sheet_names:
                print(f"  📊 Обработка листа: {sheet_name}")
                df = pd.read_excel(xlsx_path, sheet_name=sheet_name)
                
                # Конвертируем в текст
                sheet_text = f"=== ЛИСТ: {sheet_name} ===\n"
                
                # Добавляем заголовки
                if not df.empty:
                    headers = " | ".join([str(col) for col in df.columns])
                    sheet_text += f"ЗАГОЛОВКИ: {headers}\n\n"
                    
                    # Добавляем данные (ограничиваем количество строк)
                    for idx, row in df.head(1000).iterrows():
                        row_text = " | ".join([str(val) if pd.notna(val) else "" for val in row])
                        if row_text.strip():
                            sheet_text += f"{row_text}\n"
                
                all_text.append(sheet_text)
            
            combined_text = "\n\n".join(all_text)
            print(f"✅ XLSX обработан: {len(combined_text)} символов")
            return combined_text
            
        except Exception as e:
            print(f"❌ Ошибка XLSX: {e}")
            return None

    def extract_text_from_xml(self, xml_path):
        """Извлечение данных из XML файлов"""
        print(f"📄 Обработка XML: {os.path.basename(xml_path)}")
        
        try:
            tree = ET.parse(xml_path)
            root = tree.getroot()
            
            def extract_xml_text(element, level=0):
                text = []
                indent = "  " * level
                
                if element.text and element.text.strip():
                    text.append(f"{indent}{element.tag}: {element.text.strip()}")
                elif element.tag:
                    text.append(f"{indent}{element.tag}:")
                
                for child in element:
                    text.extend(extract_xml_text(child, level + 1))
                
                return text
            
            all_text = extract_xml_text(root)
            combined_text = "\n".join(all_text)
            print(f"✅ XML обработан: {len(combined_text)} символов")
            return combined_text
            
        except Exception as e:
            print(f"❌ Ошибка XML: {e}")
            return None

    def extract_text_from_pptx(self, pptx_path):
        """Извлечение текста из PowerPoint файлов"""
        print(f"📄 Обработка PPTX: {os.path.basename(pptx_path)}")
        
        try:
            from pptx import Presentation
            
            prs = Presentation(pptx_path)
            text = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = f"=== СЛАЙД {i+1} ===\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += f"{shape.text}\n"
                
                if slide_text.strip() != f"=== СЛАЙД {i+1} ===":
                    text.append(slide_text)
            
            combined_text = "\n\n".join(text)
            print(f"✅ PPTX обработан: {len(combined_text)} символов")
            return combined_text
            
        except ImportError:
            print("⚠️ python-pptx не установлен, пропускаем PPTX файл")
            return f"[PPTX файл {os.path.basename(pptx_path)} требует установки python-pptx]"
        except Exception as e:
            print(f"❌ Ошибка PPTX: {e}")
            return None

    def process_file(self, file_path):
        """Универсальная обработка файла по расширению"""
        file_ext = os.path.splitext(file_path)[1].lower()
        
        if file_ext == '.pdf':
            return self.extract_text_from_pdf(file_path)
        elif file_ext == '.docx':
            return self.extract_text_from_docx(file_path)
        elif file_ext == '.doc':
            return self.extract_text_from_doc(file_path)
        elif file_ext in ['.xlsx', '.xls']:
            return self.extract_text_from_xlsx(file_path)
        # elif file_ext == '.xml':
        #     return self.extract_text_from_xml(file_path)
        elif file_ext in ['.pptx', '.ppt']:
            return self.extract_text_from_pptx(file_path)
        else:
            print(f"⚠️ Неподдерживаемый формат: {file_ext}")
            return f"[Файл {os.path.basename(file_path)} - неподдерживаемый формат {file_ext}]"

    def process_all_files(self):
        """Обработка всех файлов в папке"""
        if not os.path.exists(self.input_dir):
            print(f"❌ Папка {self.input_dir} не найдена!")
            return False
        
        files = [f for f in os.listdir(self.input_dir) 
                if os.path.isfile(os.path.join(self.input_dir, f))]
        
        if not files:
            print(f"❌ Нет файлов в папке {self.input_dir}")
            return False
        
        print(f"🚀 Начинаем обработку {len(files)} файлов...")
        print("=" * 60)
        
        for file in files:
            file_path = os.path.join(self.input_dir, file)
            print(f"\n📁 Файл: {file}")
            
            try:
                # Обрабатываем файл
                text = self.process_file(file_path)
                
                if text:
                    # Сохраняем результат
                    output_file = os.path.join(self.output_dir, f"{file}.txt")
                    with open(output_file, 'w', encoding='utf-8') as f:
                        f.write(f"ИСТОЧНИК: {file}\n")
                        f.write(f"ДАТА ОБРАБОТКИ: {datetime.now().isoformat()}\n")
                        f.write(f"РАЗМЕР ТЕКСТА: {len(text)} символов\n")
                        f.write("=" * 60 + "\n\n")
                        f.write(text)
                    
                    self.processed_files[file] = {
                        'status': 'success',
                        'text_length': len(text),
                        'output_file': output_file
                    }
                    
                    print(f"✅ Сохранено: {output_file}")
                else:
                    self.processed_files[file] = {
                        'status': 'failed',
                        'error': 'Не удалось извлечь текст'
                    }
                    print(f"❌ Не удалось обработать файл")
                    
            except Exception as e:
                self.processed_files[file] = {
                    'status': 'error',
                    'error': str(e)
                }
                print(f"❌ Ошибка: {e}")
        
        # Создаем отчет
        self.create_processing_report()
        
        return True

    def create_processing_report(self):
        """Создаем отчет о обработке"""
        report_file = os.path.join(self.output_dir, "processing_report.txt")
        
        successful = sum(1 for f in self.processed_files.values() if f['status'] == 'success')
        failed = len(self.processed_files) - successful
        
        with open(report_file, 'w', encoding='utf-8') as f:
            f.write("ОТЧЕТ О ОБРАБОТКЕ ДОКУМЕНТОВ\n")
            f.write("=" * 40 + "\n")
            f.write(f"Дата: {datetime.now().isoformat()}\n")
            f.write(f"Всего файлов: {len(self.processed_files)}\n")
            f.write(f"Успешно обработано: {successful}\n")
            f.write(f"Ошибки: {failed}\n\n")
            
            f.write("ДЕТАЛИ ОБРАБОТКИ:\n")
            f.write("-" * 40 + "\n")
            
            for filename, info in self.processed_files.items():
                f.write(f"\nФайл: {filename}\n")
                f.write(f"Статус: {info['status']}\n")
                
                if info['status'] == 'success':
                    f.write(f"Размер текста: {info['text_length']} символов\n")
                    f.write(f"Выходной файл: {info['output_file']}\n")
                else:
                    f.write(f"Ошибка: {info.get('error', 'Неизвестная ошибка')}\n")
        
        print(f"\n📊 ИТОГИ ОБРАБОТКИ:")
        print(f"✅ Успешно: {successful}")
        print(f"❌ Ошибки: {failed}")
        print(f"📄 Отчет сохранен: {report_file}")

def main():
    processor = UniversalDocumentProcessor()
    
    print("🚀 Универсальный обработчик документов")
    print("=" * 50)
    
    # Показываем что будем обрабатывать
    if os.path.exists(processor.input_dir):
        files = os.listdir(processor.input_dir)
        print(f"📁 Найдено файлов в {processor.input_dir}: {len(files)}")
        
        # Группируем по типам
        file_types = {}
        for file in files:
            ext = os.path.splitext(file)[1].lower()
            file_types[ext] = file_types.get(ext, 0) + 1
        
        print("📊 Типы файлов:")
        for ext, count in file_types.items():
            print(f"  {ext}: {count}")
        
        print("\n🔄 Начинаем обработку...")
        processor.process_all_files()
        
    else:
        print(f"❌ Папка {processor.input_dir} не найдена!")
        print("Создайте папку и поместите туда документы для обработки")

if __name__ == "__main__":
    main()